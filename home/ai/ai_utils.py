import csv
import io
import logging
from pathlib import Path
import tempfile

import PyPDF2

logger = logging.getLogger(__name__)

try:
    from docx import Document as DocxDocument
except Exception:  # pragma: no cover - optional dependency
    DocxDocument = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency
    Presentation = None

try:
    import pytesseract
except Exception:  # pragma: no cover - optional dependency
    pytesseract = None

try:
    from PIL import Image
except Exception:  # pragma: no cover - optional dependency
    Image = None

try:
    import textract
except Exception:  # pragma: no cover - optional dependency
    textract = None


SUPPORTED_CURRICULUM_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".txt",
    ".csv",
    ".jpg",
    ".jpeg",
    ".png",
}

SUPPORTED_CURRICULUM_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "text/plain",
    "text/csv",
    "application/csv",
    "application/vnd.ms-excel",
    "image/jpeg",
    "image/png",
}


def _read_file_bytes(file_stream):
    file_stream.seek(0)
    data = file_stream.read()
    file_stream.seek(0)
    return data


def extract_text_from_pdf(pdf_file_stream):
    """Extract plain text from a PDF file-like object."""
    text = ""
    try:
        pdf_file_stream.seek(0)
        reader = PyPDF2.PdfReader(pdf_file_stream)
        for page in reader.pages:
            text += page.extract_text() or ""
        return text.strip()
    except Exception:
        logger.exception("Error extracting text from PDF")
        return None


def _extract_text_from_docx(file_stream):
    if DocxDocument is None:
        raise ValueError("DOCX extraction is unavailable on this server.")
    data = _read_file_bytes(file_stream)
    document = DocxDocument(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs if p.text).strip()


def _extract_text_from_pptx(file_stream):
    if Presentation is None:
        raise ValueError("PPTX extraction is unavailable on this server.")
    data = _read_file_bytes(file_stream)
    presentation = Presentation(io.BytesIO(data))
    lines = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines).strip()


def _extract_text_from_plain_text(file_stream):
    data = _read_file_bytes(file_stream)
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return data.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    raise ValueError("Text file encoding is not supported.")


def _extract_text_from_csv(file_stream):
    decoded = _extract_text_from_plain_text(file_stream)
    reader = csv.reader(io.StringIO(decoded))
    lines = [" | ".join(cell.strip() for cell in row if cell and cell.strip()) for row in reader]
    return "\n".join(line for line in lines if line).strip()


def _extract_text_from_image(file_stream):
    if pytesseract is None or Image is None:
        raise ValueError("Image OCR extraction is unavailable on this server.")
    data = _read_file_bytes(file_stream)
    image = Image.open(io.BytesIO(data))
    return (pytesseract.image_to_string(image) or "").strip()


def _extract_with_textract(file_stream, filename):
    if textract is None:
        raise ValueError("This file type is unsupported on this server.")
    data = _read_file_bytes(file_stream)
    suffix = Path(filename or "").suffix.lower()
    if not suffix or suffix not in SUPPORTED_CURRICULUM_EXTENSIONS:
        suffix = ".bin"
    with tempfile.NamedTemporaryFile(suffix=suffix) as temp_file:
        temp_file.write(data)
        temp_file.flush()
        extracted = textract.process(temp_file.name)
    return extracted.decode("utf-8", errors="ignore").strip()


def extract_text_from_file(file_stream, filename, content_type=None):
    """Extract text from a supported uploaded curriculum file."""
    ext = Path(filename or "").suffix.lower()
    logger.info("Attempting extraction for file '%s' (ext=%s, type=%s)", filename, ext, content_type)

    try:
        if ext == ".pdf":
            text = extract_text_from_pdf(file_stream)
        elif ext == ".docx":
            text = _extract_text_from_docx(file_stream)
        elif ext == ".pptx":
            text = _extract_text_from_pptx(file_stream)
        elif ext == ".txt":
            text = _extract_text_from_plain_text(file_stream)
        elif ext == ".csv":
            text = _extract_text_from_csv(file_stream)
        elif ext in {".jpg", ".jpeg", ".png"}:
            text = _extract_text_from_image(file_stream)
        else:
            text = _extract_with_textract(file_stream, filename)

        if not text:
            raise ValueError("No extractable text was found in this file.")

        logger.info("Extraction successful for file '%s'", filename)
        return text
    except Exception as exc:
        logger.exception("Extraction failed for file '%s': %s", filename, exc)
        if isinstance(exc, ValueError):
            raise
        raise ValueError("Unable to process this document. Please upload a supported, non-corrupted file.") from exc
